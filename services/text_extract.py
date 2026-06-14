"""
Plain-text extraction from uploaded documents.

Replaces the old services/extractor.py + ChromaDB ingestion pipeline.
One function, one job: take a file path, return a string. The caller
stashes that string with the chat session so it can be cached with
Gemini's context-cache API.
"""

import logging
from pathlib import Path

logger = logging.getLogger(__name__)

SUPPORTED_EXTS = {".pdf", ".docx", ".pptx", ".txt", ".md"}

# Hard cap to keep us under Gemini 2.5 Flash's 1M-token input limit with headroom.
# ~3.5 chars/token average; 900_000 chars ≈ 257K tokens. Truncate beyond this.
MAX_CHARS = 900_000


def extract_text(file_path: str) -> str:
    """
    Extract plain UTF-8 text from a document. Returns '' if the file
    has no extractable text (e.g. image-only scanned PDF) so callers
    can warn the user instead of crashing.

    Raises ValueError for unsupported extensions.
    """
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        text = _from_pdf(file_path)
    elif ext == ".docx":
        text = _from_docx(file_path)
    elif ext == ".pptx":
        text = _from_pptx(file_path)
    elif ext in (".txt", ".md"):
        text = _from_text(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    if len(text) > MAX_CHARS:
        logger.warning(
            "[text_extract] %s truncated from %d to %d chars",
            Path(file_path).name, len(text), MAX_CHARS,
        )
        text = text[:MAX_CHARS]
    return text


def _from_pdf(path: str) -> str:
    try:
        import fitz  # PyMuPDF
    except ImportError:
        logger.error("[text_extract] PyMuPDF (fitz) not installed")
        return ""
    parts = []
    with fitz.open(path) as doc:
        for page in doc:
            t = page.get_text().strip()
            if t:
                parts.append(t)
    return "\n\n".join(parts)


def _from_docx(path: str) -> str:
    try:
        import docx
    except ImportError:
        logger.error("[text_extract] python-docx not installed")
        return ""
    doc = docx.Document(path)
    parts = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    return "\n".join(parts)


def _from_pptx(path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        logger.error("[text_extract] python-pptx not installed")
        return ""
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        bits = []
        title_shape = slide.shapes.title if slide.shapes.title else None
        if title_shape and title_shape.text and title_shape.text.strip():
            bits.append(title_shape.text.strip())
        for shape in slide.shapes:
            if shape is title_shape:
                continue
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                bits.append(shape.text.strip())
        if bits:
            slides.append(f"--- Slide {i + 1} ---\n" + "\n".join(bits))
    return "\n\n".join(slides)


def _from_text(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()
