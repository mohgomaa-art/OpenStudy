import os
from pathlib import Path
from typing import List, Dict, Any

class Node:
    def __init__(self, node_id: str, title: str, content: str, parent_id: str = None):
        self.node_id = node_id
        self.title = title
        self.content = content
        self.parent_id = parent_id
        self.children = []

    def to_dict(self):
        return {
            "node_id": self.node_id,
            "title": self.title,
            "content": self.content,
            "parent_id": self.parent_id
        }

def extract_structure(file_path: str) -> List[Node]:
    """Route extraction based on file extension."""
    ext = Path(file_path).suffix.lower()
    if ext == ".pptx":
        return _extract_pptx(file_path)
    elif ext == ".pdf":
        return _extract_pdf(file_path)
    elif ext == ".docx":
        return _extract_docx(file_path)
    elif ext == ".txt":
        return _extract_txt(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

def _extract_pptx(file_path: str) -> List[Node]:
    try:
        from pptx import Presentation
    except ImportError:
        print("[Extractor] python-pptx not installed")
        return []
    
    prs = Presentation(file_path)
    nodes = []
    
    # Root node for the presentation
    root_id = f"doc_{Path(file_path).stem}"
    root = Node(root_id, Path(file_path).name, "Presentation Root")
    nodes.append(root)

    for i, slide in enumerate(prs.slides):
        title = f"Slide {i+1}"
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()
            
        content_lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip() and shape != slide.shapes.title:
                content_lines.append(shape.text.strip())
                
        content = "\n".join(content_lines)
        if not content:
            content = title
            
        slide_node = Node(f"{root_id}_slide_{i}", title, content, root_id)
        nodes.append(slide_node)
        
    return nodes

def _extract_pdf(file_path: str) -> List[Node]:
    try:
        import fitz # PyMuPDF
    except ImportError:
        print("[Extractor] PyMuPDF not installed")
        return []
        
    doc = fitz.open(file_path)
    nodes = []
    root_id = f"doc_{Path(file_path).stem}"
    nodes.append(Node(root_id, Path(file_path).name, "PDF Root"))
    
    # Simple extraction: each page is a node (for true heading hierarchy, we'd need TOC/Bookmarks)
    toc = doc.get_toc()
    if toc:
        # Use TOC for hierarchy if available
        for i, item in enumerate(toc):
            lvl, title, page = item
            try:
                content = doc[page-1].get_text().strip()
            except Exception:
                content = ""
            nodes.append(Node(f"{root_id}_toc_{i}", title, content, root_id))
    else:
        # Fallback: pages
        for i in range(len(doc)):
            text = doc[i].get_text().strip()
            if text:
                title = f"Page {i+1}"
                nodes.append(Node(f"{root_id}_page_{i}", title, text, root_id))
                
    return nodes

def _extract_docx(file_path: str) -> List[Node]:
    try:
        import docx
    except ImportError:
        print("[Extractor] python-docx not installed")
        return []
        
    doc = docx.Document(file_path)
    nodes = []
    root_id = f"doc_{Path(file_path).stem}"
    nodes.append(Node(root_id, Path(file_path).name, "DOCX Root"))
    
    current_node = None
    current_content = []
    node_idx = 0
    
    for para in doc.paragraphs:
        style = para.style.name
        text = para.text.strip()
        if not text:
            continue
            
        if style.startswith('Heading'):
            if current_node:
                current_node.content = "\n".join(current_content)
                nodes.append(current_node)
                
            current_node = Node(f"{root_id}_h_{node_idx}", text, "", root_id)
            current_content = []
            node_idx += 1
        else:
            if not current_node:
                current_node = Node(f"{root_id}_intro", "Introduction", "", root_id)
            current_content.append(text)
            
    if current_node:
        current_node.content = "\n".join(current_content)
        nodes.append(current_node)
        
    return nodes

def _extract_txt(file_path: str) -> List[Node]:
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        content = f.read()
        
    nodes = []
    root_id = f"doc_{Path(file_path).stem}"
    nodes.append(Node(root_id, Path(file_path).name, "Text Root"))
    
    paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
    for i, p in enumerate(paragraphs):
        title = " ".join(p.split()[:5]) + "..."
        nodes.append(Node(f"{root_id}_p_{i}", title, p, root_id))
        
    return nodes
