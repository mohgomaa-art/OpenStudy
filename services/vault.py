import os
import json
import hashlib
import asyncio
import tempfile
from pathlib import Path
from typing import AsyncGenerator
from dotenv import load_dotenv

from services.config import config_service

# VAULT_PATH و MANIFEST_PATH يتم ضبطهم بناءً على root_path الموحد
VAULT_PATH   = os.getenv("VAULT_PATH") or str(config_service.data_path / "data" / "vault")
MANIFEST_PATH = config_service.data_path / "data" / "vault_manifest.json"
SUPPORTED_EXT = {
    ".pdf", ".docx", ".txt", ".md",
    ".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".svg",
    ".mp3", ".wav", ".m4a", ".ogg", ".flac", ".aac",
    ".mp4", ".mkv", ".mov", ".avi", ".webm",
    ".pptx", ".xlsx",
    ".ytlink", ".weblink", ".flashcard", ".session",
}


# ── Text Extractors ───────────────────────────────────────────────────────────

_PDF_TEXT_FALLBACK_THRESHOLD = 80


def _extract_pdf_text_from_doc(doc, *, include_page_markers: bool = False) -> str:
    parts = []
    for i, page in enumerate(doc):
        text = page.get_text("text").strip()
        if not text:
            continue
        if include_page_markers:
            parts.append(f"=== Page {i + 1} ===\n{text}")
        else:
            parts.append(text)
    return "\n\n".join(parts).strip()


def _ocr_pdf_from_doc(doc, *, include_page_markers: bool = False) -> str:
    try:
        import fitz  # pymupdf
        from PIL import Image
        import pytesseract

        scale = max(2.0, float(os.getenv("PDF_OCR_SCALE", "2.2")))
        matrix = fitz.Matrix(scale, scale)
        lang = os.getenv("TESSERACT_LANGS", "ara+eng")
        parts = []
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            image = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            text = pytesseract.image_to_string(image, lang=lang).strip()
            if not text:
                continue
            if include_page_markers:
                parts.append(f"=== Page {i + 1} ===\n{text}")
            else:
                parts.append(text)
        return "\n\n".join(parts).strip()
    except Exception as e:
        return f"[PDF OCR Error: {e}]"


def extract_pdf_bytes(data: bytes, *, force_ocr: bool = False, include_page_markers: bool = False) -> str:
    try:
        import fitz  # pymupdf

        doc = fitz.open(stream=data, filetype="pdf")
        try:
            embedded_text = _extract_pdf_text_from_doc(doc, include_page_markers=include_page_markers)
            if not force_ocr and len(embedded_text) >= _PDF_TEXT_FALLBACK_THRESHOLD:
                return embedded_text

            ocr_text = _ocr_pdf_from_doc(doc, include_page_markers=include_page_markers)
            if force_ocr:
                return ocr_text if ocr_text and not ocr_text.startswith("[") else embedded_text or ocr_text

            if ocr_text and not ocr_text.startswith("[") and len(ocr_text) > len(embedded_text):
                return ocr_text
            return embedded_text or ocr_text
        finally:
            doc.close()
    except Exception as e:
        return f"[PDF Error: {e}]"


def _extract_pdf(path: str, *, force_ocr: bool = False) -> str:
    """PyMuPDF text extraction with OCR fallback for scanned PDFs."""
    try:
        return extract_pdf_bytes(Path(path).read_bytes(), force_ocr=force_ocr)
    except Exception as e:
        return f"[PDF Error: {e}]"


def _extract_docx(path: str) -> str:
    try:
        import mammoth
        with open(path, "rb") as f:
            result = mammoth.extract_raw_text(f)
        return result.value.strip()
    except Exception as e:
        return f"[DOCX Error: {e}]"


def _extract_text(path: str) -> str:
    try:
        with open(path, encoding="utf-8", errors="ignore") as f:
            return f.read().strip()
    except Exception as e:
        return f"[Text Error: {e}]"


def extract_text(path: str, *, force_ocr: bool = False) -> str:
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(path, force_ocr=force_ocr)
    elif ext == ".docx":
        return _extract_docx(path)
    elif ext in {".txt", ".md"}:
        return _extract_text(path)
    elif ext in {".png", ".jpg", ".jpeg"}:
        return _extract_image(path)
    elif ext == ".pptx":
        return _extract_pptx(path)
    elif ext == ".xlsx":
        return _extract_xlsx(path)
    return ""


def _extract_image(path: str) -> str:
    try:
        from PIL import Image
        import pytesseract
        img = Image.open(path)
        text = pytesseract.image_to_string(img, lang="ara+eng")
        return text.strip()
    except Exception as e:
        return f"[Image OCR Error: {e}]"


def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    parts.append(shape.text)
        return "\n".join(p for p in parts if p.strip()).strip()
    except Exception as e:
        return f"[PPTX Error: {e}]"


def _extract_xlsx(path: str) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True)
        parts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                row_text = " ".join(str(c) for c in row if c is not None).strip()
                if row_text:
                    parts.append(row_text)
        return "\n".join(parts).strip()
    except Exception as e:
        return f"[XLSX Error: {e}]"


def _file_hash(path: str) -> str:
    """SHA256 لأول 64KB من الملف — للكشف عن التغييرات."""
    h = hashlib.sha256()
    with open(path, "rb") as f:
        h.update(f.read(65536))
    return h.hexdigest()[:16]


async def _extract_audio_text(path: str, transcribe_fn) -> str:
    """Extract text from audio/video files using local STT."""
    try:
        from pydub import AudioSegment
        audio = AudioSegment.from_file(path)
        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp:
            audio.export(tmp.name, format="wav")
            wav_path = tmp.name
        try:
            audio_bytes = Path(wav_path).read_bytes()
            result = await transcribe_fn(audio_bytes, language="ar")
            return result.get("text", "").strip()
        finally:
            Path(wav_path).unlink(missing_ok=True)
    except Exception as e:
        return f"[Audio Error: {e}]"


# ── Manifest ──────────────────────────────────────────────────────────────────

def _load_manifest() -> dict:
    if MANIFEST_PATH.exists():
        try:
            return json.loads(MANIFEST_PATH.read_text())
        except Exception:
            pass
    return {}


def _save_manifest(manifest: dict):
    MANIFEST_PATH.parent.mkdir(parents=True, exist_ok=True)
    MANIFEST_PATH.write_text(json.dumps(manifest, ensure_ascii=False, indent=2))


# ── Main Indexer ──────────────────────────────────────────────────────────────

class VaultIndexer:

    def __init__(self):
        self._manifest = _load_manifest()
        self._indexed_files: list[str] = list(self._manifest.keys())
        self._tree_cache = None
        self._cache_lock = asyncio.Lock()
        self._llm_semaphore = asyncio.Semaphore(3)  # أقصى 3 LLM calls في نفس الوقت

    async def index(self, vault_path: str) -> AsyncGenerator[str, None]:
        """
        يفهرس الـ vault ويبثّ progress عبر SSE.
        يتخطى الملفات اللي لم تتغير (بناءً على الـ hash).
        """
        from services.rag import add_document
        from services.config import config_service
        from services.audio import transcribe_audio
        from services.brain_map import extract_prerequisites
        from services.memory import memory_service

        if not vault_path or not vault_path.strip():
            vault_path = config_service.ensure_vault_path()
            print(f"[VaultIndexer] Empty path provided, fallback to: {vault_path}")

        p = Path(vault_path).resolve()
        
        # Prevent indexing the project root or its parents (safety)
        project_root = config_service.root_path
        if p == project_root or project_root in p.parents:
             if not (p.name == "vault" and p.parent.name == "data"):
                print(f"[VaultIndexer] Warning: Indexing project-related path: {p}")

        if not p.exists():
            fallback = Path(config_service.ensure_vault_path()).resolve()
            if fallback.exists():
                p = fallback
                vault_path = str(p)
            else:
                yield f"data: {json.dumps({'error': 'Path not found'})}\n\n"
                return

        if p.is_file():
            all_files = [str(p)] if p.suffix.lower() in SUPPORTED_EXT else []
        else:
            all_files = [
                str(f) for f in p.rglob("*")
                if f.is_file() and f.suffix.lower() in SUPPORTED_EXT
            ]
        
        total = len(all_files)
        if not all_files:
            yield f"data: {json.dumps({'percent': 100, 'file': 'No files found', 'total': 0, 'empty': True, 'vault_path': vault_path})}\n\n"
            return

        for i, file_path in enumerate(all_files):
            file_hash = _file_hash(file_path)
            file_name = Path(file_path).name

            if self._manifest.get(file_path) == file_hash:
                update = {"percent": int((i + 1) / total * 100), "file": file_name, "skipped": True, "total": total}
                yield f"data: {json.dumps(update)}\n\n"
                await asyncio.sleep(0)
                continue

            ext = Path(file_path).suffix.lower()
            source_type = "image" if ext in {".png", ".jpg", ".jpeg", ".webp", ".gif"} else "document"
            
            if ext in {".mp3", ".wav", ".mp4"}:
                text = await _extract_audio_text(file_path, transcribe_audio)
                source_type = "audio"
            else:
                # BUG-05 FIX: run CPU-bound extraction off the event loop
                text = await asyncio.to_thread(extract_text, file_path)

            if text and not text.startswith("["):
                chunks_added = add_document(
                    text=text,
                    source=file_path,
                    metadata={"filename": file_name, "hash": file_hash, "source_type": source_type},
                )
                
                # Full-Text Search Indexing
                memory_service.update_vault_index(file_path, file_name, text, source_type=source_type)
                
                # Semantic Extraction — WASTE-01 FIX: only for newly indexed text files
                # Guard against re-running on every index pass (costs ~900 tokens/file)
                semantic_key = f"semantic:{file_path}"
                if (
                    ext in {".pdf", ".docx", ".md", ".txt"}
                    and len(text) > 500
                    and not self._manifest.get(semantic_key)
                ):
                    try:
                        async with self._llm_semaphore:
                            await extract_prerequisites(text, file_path)
                        self._manifest[semantic_key] = True
                    except Exception as e:
                        print(f"[VaultIndexer] Semantic Extraction Failed: {e}")

                self._manifest[file_path] = file_hash
                if file_path not in self._indexed_files:
                    self._indexed_files.append(file_path)

                update = {
                    "percent": int((i + 1) / total * 100),
                    "file": file_name,
                    "chunks": chunks_added,
                    "total": total,
                }
            else:
                update = {
                    "percent": int((i + 1) / total * 100),
                    "file": file_name,
                    "error": "Text extraction failed",
                    "total": total,
                }

            yield f"data: {json.dumps(update)}\n\n"
            await asyncio.sleep(0.01)

        dead_files = [f for f in self._manifest.keys() if f not in all_files]
        for f in dead_files:
            del self._manifest[f]
            if f in self._indexed_files:
                self._indexed_files.remove(f)

        _save_manifest(self._manifest)
        async with self._cache_lock:
            self._tree_cache = None # Invalidate cache after index
        yield f"data: {json.dumps({'percent': 100, 'file': 'Indexing complete!', 'done': True, 'total': total})}\n\n"

    async def get_tree(self) -> list:
        """Returns a cached tree of the vault files."""
        from services.config import config_service
        
        async with self._cache_lock:
            if self._tree_cache is not None:
                return self._tree_cache

            vault_path = Path(config_service.ensure_vault_path()).resolve()
            items = []
            if vault_path.exists():
                # 1. Fetch metadata lookup from memory.db
                from services.memory import memory_service
                meta_lookup = {}
                with memory_service._lock:
                    cur = memory_service._conn.execute("SELECT path, tags, subject FROM file_metadata")
                    for row in cur.fetchall():
                        meta_lookup[row[0]] = {
                            "tags": json.loads(row[1] or "[]"),
                            "subject": row[2]
                        }

                # 2. Run CPU-bound walk in a thread
                def _walk():
                    res = []
                    for root, dirs, files in os.walk(vault_path):
                        p_root = Path(root)
                        if p_root.name.startswith(".") or p_root.name == "_quick_notes" or ".git" in p_root.parts:
                            continue
                        try:
                            rel_dir = p_root.relative_to(vault_path).as_posix()
                        except ValueError:
                            continue
                        if rel_dir == ".": rel_dir = ""
                        for d in dirs:
                            if d.startswith(".") or d == "_quick_notes": continue
                            d_path = str((p_root / d).resolve())
                            meta = meta_lookup.get(d_path, {})
                            res.append({
                                "name": d, 
                                "type": "folder", 
                                "path": d_path,
                                "parent": rel_dir,
                                "tags": meta.get("tags", []),
                                "subject": meta.get("subject")
                            })
                        for f in files:
                            if f.startswith("."): continue
                            f_path = p_root / f
                            if f_path.suffix.lower() in SUPPORTED_EXT:
                                full_path = str(f_path.resolve())
                                meta = meta_lookup.get(full_path, {})
                                res.append({
                                    "name": f,
                                    "type": "file",
                                    "path": full_path,
                                    "parent": rel_dir,
                                    "tags": meta.get("tags", []),
                                    "subject": meta.get("subject")
                                })
                    return res
                
                items = await asyncio.to_thread(_walk)
            
            self._tree_cache = items
            return items

    def invalidate_cache(self):
        self._tree_cache = None


    def get_stats(self) -> dict:
        from services.memory import memory_service
        from services.rag import get_chunk_count

        return {
            "file_count":   len(self._indexed_files),
            "chunk_count":  get_chunk_count(),
            "readiness":    memory_service.readiness_score(),
            "weak_topics":  memory_service.get_weak_topics(limit=5),
            "weekly_stats": memory_service.get_weekly_stats(),
            "last_studied": self._last_activity(),
        }

    def _last_activity(self) -> str:
        if not self._indexed_files:
            return "No activity yet"
        try:
            mtime = max(os.path.getmtime(f) for f in self._indexed_files if os.path.exists(f))
            from datetime import datetime
            diff = datetime.now() - datetime.fromtimestamp(mtime)
            if diff.days == 0:
                hours = int(diff.seconds / 3600)
                return f"{hours}h ago" if hours > 0 else "Just now"
            return f"{diff.days}d ago"
        except Exception:
            return "Unknown"

    @property
    def indexed_files(self) -> list[str]:
        return self._indexed_files


indexer = VaultIndexer()
